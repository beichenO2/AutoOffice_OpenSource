#!/usr/bin/env python3
"""从 stdin 读入 JSON，写出 argv[1] 路径的 .pptx（python-pptx）。"""
from __future__ import annotations

import base64
from io import BytesIO
import json
import sys
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

MAX_BULLETS = 6

# 多主题：配色、字号与左侧装饰条（minimal 无条，偏「留白杂志」感）
THEMES: dict[str, dict[str, Any]] = {
    "academic": {
        "title_rgb": (0x1A, 0x36, 0x5D),
        "body_rgb": (0x2D, 0x37, 0x41),
        "muted_rgb": (0x4A, 0x55, 0x66),
        "accent_rgb": (0x2C, 0x5A, 0x8F),
        "accent_bar": True,
        "accent_width_in": 0.1,
        "title_pt": 32,
        "bullet_pt": 17,
        "title_font": "PingFang SC",
        "body_font": "PingFang SC",
    },
    "business": {
        "title_rgb": (0x0F, 0x17, 0x2A),
        "body_rgb": (0x1E, 0x29, 0x3B),
        "muted_rgb": (0x5C, 0x6B, 0x7A),
        "accent_rgb": (0x30, 0x58, 0x8C),
        "accent_bar": True,
        "accent_width_in": 0.12,
        "title_pt": 34,
        "bullet_pt": 18,
        "title_font": "Helvetica Neue",
        "body_font": "PingFang SC",
    },
    "minimal": {
        "title_rgb": (0x14, 0x14, 0x14),
        "body_rgb": (0x33, 0x33, 0x33),
        "muted_rgb": (0x66, 0x66, 0x66),
        "accent_bar": False,
        "title_pt": 36,
        "bullet_pt": 20,
        "title_font": "PingFang SC",
        "body_font": "PingFang SC",
    },
    "nord": {
        "title_rgb": (0x2E, 0x34, 0x40),
        "body_rgb": (0x3B, 0x42, 0x52),
        "muted_rgb": (0x4C, 0x56, 0x6A),
        "accent_rgb": (0x88, 0xC0, 0xD0),
        "accent_bar": True,
        "accent_width_in": 0.11,
        "title_pt": 33,
        "bullet_pt": 18,
        "title_font": "PingFang SC",
        "body_font": "PingFang SC",
    },
    "tech": {
        "title_rgb": (0x0F, 0x17, 0x2A),
        "body_rgb": (0x1E, 0x29, 0x3B),
        "muted_rgb": (0x64, 0x74, 0x8B),
        "accent_rgb": (0x38, 0xBD, 0xF8),
        "accent_bar": True,
        "accent_width_in": 0.08,
        "title_pt": 33,
        "bullet_pt": 17,
        "title_font": "Helvetica Neue",
        "body_font": "PingFang SC",
    },
    "warm": {
        "title_rgb": (0x5C, 0x40, 0x33),
        "body_rgb": (0x3D, 0x2F, 0x28),
        "muted_rgb": (0x6B, 0x5B, 0x50),
        "accent_rgb": (0xC4, 0x8B, 0x6A),
        "accent_bar": True,
        "accent_width_in": 0.13,
        "title_pt": 32,
        "bullet_pt": 18,
        "title_font": "PingFang SC",
        "body_font": "PingFang SC",
    },
    "slate": {
        "title_rgb": (0x33, 0x41, 0x55),
        "body_rgb": (0x47, 0x55, 0x69),
        "muted_rgb": (0x71, 0x85, 0xA0),
        "accent_rgb": (0x63, 0x66, 0xF1),
        "accent_bar": True,
        "accent_width_in": 0.1,
        "title_pt": 34,
        "bullet_pt": 18,
        "title_font": "PingFang SC",
        "body_font": "PingFang SC",
    },
}


def _rgb(t: tuple[int, int, int]) -> RGBColor:
    return RGBColor(t[0], t[1], t[2])


def _add_left_accent(slide: Any, prs: Presentation, th: dict[str, Any]) -> None:
    if not th.get("accent_bar"):
        return
    w_in = float(th.get("accent_width_in", 0.1))
    accent = th.get("accent_rgb", th["title_rgb"])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        Inches(w_in),
        prs.slide_height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(accent)  # type: ignore[arg-type]
    if shape.line:
        shape.line.fill.background()


def _style_title(shape: Any, th: dict[str, Any]) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for run in p.runs:
            run.font.name = th["title_font"]
            run.font.size = Pt(th["title_pt"])
            run.font.bold = True
            run.font.color.rgb = _rgb(th["title_rgb"])


def _style_body(tf: Any, th: dict[str, Any]) -> None:
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.space_after = Pt(6)
        for run in p.runs:
            run.font.name = th["body_font"]
            run.font.size = Pt(th["bullet_pt"])
            run.font.bold = False
            run.font.color.rgb = _rgb(th["body_rgb"])


def build(data: dict[str, Any], out_path: str) -> None:
    theme_id = data.get("theme", "minimal")
    th = THEMES.get(theme_id, THEMES["minimal"])

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # 封面
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])
    _add_left_accent(slide0, prs, th)
    slide0.shapes.title.text = data.get("title", "")
    _style_title(slide0.shapes.title, th)
    sub = data.get("subtitle") or ""
    if len(slide0.placeholders) > 1:
        ph_sub = slide0.placeholders[1]
        ph_sub.text = sub
        _style_subtitle(ph_sub, th)

    for s in data.get("slides", []):
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        _add_left_accent(sl, prs, th)
        sl.shapes.title.text = s.get("title", "")
        _style_title(sl.shapes.title, th)

        bullets = (s.get("bullets") or [])[:MAX_BULLETS]
        body = sl.shapes.placeholders[1].text_frame
        chart_png = s.get("chart_png_base64")
        if not bullets:
            body.text = ""
        else:
            body.text = bullets[0]
            for b in bullets[1:]:
                p = body.add_paragraph()
                p.text = b
                p.level = 0
        _style_body(body, th)

        if chart_png:
            image_stream = BytesIO(base64.b64decode(chart_png))
            sl.shapes.add_picture(image_stream, Inches(1.0), Inches(1.8), width=Inches(10.8))

    prs.save(out_path)


def _style_subtitle(shape: Any, th: dict[str, Any]) -> None:
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = th["body_font"]
            run.font.size = Pt(th["bullet_pt"] - 1)
            run.font.color.rgb = _rgb(th["muted_rgb"])


def main() -> None:
    if len(sys.argv) < 2:
        print("usage: build_pptx.py <output.pptx> < stdin.json", file=sys.stderr)
        sys.exit(2)
    raw = sys.stdin.read()
    data = json.loads(raw)
    build(data, sys.argv[1])


if __name__ == "__main__":
    main()
